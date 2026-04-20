"""
GPTHub Presentation Builder
============================
Extracts the PPTX generation pipeline from middleware.py.
Themes, layouts, slide image generation, and the chat handler live here.
"""

from __future__ import annotations

import io
import json
import logging
import re
import uuid as _uuid
from typing import Any

from fastapi import Request

from open_webui.utils.gpthub_i18n import t

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Themes
# ---------------------------------------------------------------------------

PPTX_THEMES: dict[str, dict] = {
    'corporate': {
        'bg': (0x1A, 0x2E, 0x4A),
        'title_bg': (0x12, 0x1F, 0x33),
        'title_fg': (0xFF, 0xFF, 0xFF),
        'body_fg': (0xCC, 0xDD, 0xEE),
        'html_bg': '#1a2e4a',
        'html_card': '#12213a',
        'html_title': '#ffffff',
        'html_body': '#ccddee',
        'html_accent': '#008bd0',
    },
    'dark': {
        'bg': (0x1A, 0x1A, 0x2E),
        'title_bg': (0x0D, 0x0D, 0x1A),
        'title_fg': (0xE0, 0xFA, 0xFF),
        'body_fg': (0xB0, 0xC4, 0xDE),
        'html_bg': '#1a1a2e',
        'html_card': '#0d0d1a',
        'html_title': '#e0faff',
        'html_body': '#b0c4de',
        'html_accent': '#7f5af0',
    },
    'minimal': {
        'bg': (0xFF, 0xFF, 0xFF),
        'title_bg': (0xF5, 0xF5, 0xF5),
        'title_fg': (0x1A, 0x1A, 0x2E),
        'body_fg': (0x44, 0x44, 0x55),
        'html_bg': '#f5f5f5',
        'html_card': '#ffffff',
        'html_title': '#1a1a2e',
        'html_body': '#444455',
        'html_accent': '#3366ff',
    },
}

PPTX_THEME_KEYWORDS: list[tuple[list[str], str]] = [
    (['corporate', 'business', 'professional', 'корпорат', 'бизнес', 'деловой'], 'corporate'),
    (['dark', 'night', 'темн', 'ночн', 'черн'], 'dark'),
    (['minimal', 'clean', 'light', 'white', 'минимал', 'чист', 'светл', 'белый'], 'minimal'),
]


def detect_pptx_theme(text: str) -> str:
    lowered = text.lower()
    for keywords, theme in PPTX_THEME_KEYWORDS:
        if any(kw in lowered for kw in keywords):
            return theme
    return 'corporate'


# ---------------------------------------------------------------------------
# Drain streamed/non-streamed LLM response to plain text
# ---------------------------------------------------------------------------


async def drain_response_to_text(result) -> str:
    if isinstance(result, dict):
        return (result.get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
    if hasattr(result, 'body') and not hasattr(result, 'body_iterator'):
        try:
            parsed = json.loads(result.body)
            return (parsed.get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
        except Exception:
            return ''
    if hasattr(result, 'body_iterator'):
        parts: list[str] = []
        async for chunk in result.body_iterator:
            text = chunk.decode('utf-8') if isinstance(chunk, bytes) else chunk
            for line in text.split('\n'):
                line = line.strip()
                if not line.startswith('data:'):
                    continue
                data_str = line[5:].strip()
                if not data_str or data_str == '[DONE]':
                    continue
                try:
                    data = json.loads(data_str)
                    choice = (data.get('choices') or [{}])[0]
                    piece = (choice.get('delta') or {}).get('content') or ''
                    if not piece:
                        piece = (choice.get('message') or {}).get('content') or ''
                    parts.append(piece)
                except Exception:
                    pass
        return ''.join(parts)
    raise ValueError(f'Unexpected response type from generate_chat_completion: {type(result)}')


# ---------------------------------------------------------------------------
# Slide layouts  (slide = 13.33 × 7.5 in)
# ---------------------------------------------------------------------------

_SLIDE_W = 13.33
_SLIDE_H = 7.5

PPTX_LAYOUTS: dict[str, dict[str, tuple]] = {
    'split-right': {
        'title': (0.5, 0.3, 6.0, 1.0),
        'body': (0.5, 1.5, 6.0, 5.5),
        'image': (6.83, 0.0, 6.5, 7.5),
    },
    'split-left': {
        'title': (6.83, 0.3, 6.0, 1.0),
        'body': (6.83, 1.5, 6.0, 5.5),
        'image': (0.0, 0.0, 6.5, 7.5),
    },
    'hero': {
        'title': (1.0, 2.0, 11.33, 1.5),
        'body': (1.5, 4.0, 10.33, 2.5),
        'image': (0.0, 0.0, _SLIDE_W, _SLIDE_H),
    },
    'top-strip': {
        'title': (0.5, 3.6, 12.33, 1.0),
        'body': (0.5, 4.8, 12.33, 2.4),
        'image': (0.0, 0.0, _SLIDE_W, 3.4),
    },
}

VALID_LAYOUTS = list(PPTX_LAYOUTS.keys())


# ---------------------------------------------------------------------------
# PPTX builder (pure, no IO)
# ---------------------------------------------------------------------------


def build_pptx(
    slides_data: dict,
    theme: str = 'corporate',
    slide_images: dict[int, bytes] | None = None,
) -> bytes:
    from pptx import Presentation as PPTXPresentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt

    slide_images = slide_images or {}
    t = PPTX_THEMES.get(theme, PPTX_THEMES['corporate'])

    prs = PPTXPresentation()
    prs.slide_width = Inches(_SLIDE_W)
    prs.slide_height = Inches(_SLIDE_H)
    blank_layout = prs.slide_layouts[6]

    slides = slides_data.get('slides', [])
    for idx, sd in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)

        slide.background.fill.solid()
        bg_color = t['title_bg'] if idx == 0 else t['bg']
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

        layout_name = sd.get('layout', 'split-right')
        if layout_name not in PPTX_LAYOUTS:
            layout_name = 'split-right'
        layout = PPTX_LAYOUTS[layout_name]

        # Image
        img_bytes = slide_images.get(idx)
        if img_bytes:
            il, it, iw, ih = layout['image']
            if layout_name == 'hero':
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    Inches(il), Inches(it), Inches(iw), Inches(ih),
                )
                from pptx.oxml.ns import qn
                from lxml import etree

                overlay = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(_SLIDE_W), Inches(_SLIDE_H),
                )
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(*bg_color)
                sp_pr = overlay._element.spPr
                solid = sp_pr.find(qn('a:solidFill'))
                if solid is not None:
                    srgb = solid.find(qn('a:srgbClr'))
                    if srgb is not None:
                        alpha = etree.SubElement(srgb, qn('a:alpha'))
                        alpha.set('val', '50000')
                overlay.line.fill.background()
            else:
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    Inches(il), Inches(it), Inches(iw), Inches(ih),
                )

        # Title
        tl, tt, tw, th = layout['title']
        title_box = slide.shapes.add_textbox(Inches(tl), Inches(tt), Inches(tw), Inches(th))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = sd.get('title', f'Slide {idx + 1}')
        title_p.font.size = Pt(32) if layout_name == 'hero' else Pt(24)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(*t['title_fg'])
        if layout_name == 'hero':
            title_p.alignment = PP_ALIGN.CENTER

        # Body
        body_items = sd.get('body', [])
        if body_items:
            bl, bt, bw, bh = layout['body']
            body_box = slide.shapes.add_textbox(Inches(bl), Inches(bt), Inches(bw), Inches(bh))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for j, bullet in enumerate(body_items):
                p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(16) if layout_name == 'hero' else Pt(14)
                p.font.color.rgb = RGBColor(*t['body_fg'])
                p.space_after = Pt(6)
                if layout_name == 'hero':
                    p.alignment = PP_ALIGN.CENTER

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Slide image generation
# ---------------------------------------------------------------------------


async def generate_slide_images(
    request: Request,
    slides: list[dict],
    metadata: dict,
    user,
    event_emitter,
    *,
    image_generations_fn=None,
    create_image_form_cls=None,
) -> tuple[dict[int, bytes], dict[int, str]]:
    from open_webui.models.files import Files
    from open_webui.storage.provider import Storage

    slide_images: dict[int, bytes] = {}
    slide_image_urls: dict[int, str] = {}
    total = sum(1 for s in slides if s.get('image_prompt'))
    generated = 0

    if not image_generations_fn or not create_image_form_cls:
        return slide_images, slide_image_urls

    for idx, sd in enumerate(slides):
        prompt = sd.get('image_prompt')
        if not prompt:
            continue

        generated += 1
        await event_emitter({
            'type': 'status',
            'data': {'description': t('pres.status_image_n', n=generated, total=total), 'done': False},
        })

        try:
            images = await image_generations_fn(
                request=request,
                form_data=create_image_form_cls(**{'prompt': prompt, 'size': '1024x1024', 'n': 1}),
                metadata={
                    'chat_id': metadata.get('chat_id'),
                    'message_id': metadata.get('message_id'),
                },
                user=user,
            )

            if images and images[0].get('url'):
                img_url = images[0]['url']
                slide_image_urls[idx] = img_url

                parts = img_url.strip('/').split('/')
                file_id = None
                for pi, part in enumerate(parts):
                    if part == 'files' and pi + 1 < len(parts):
                        file_id = parts[pi + 1]
                        break

                if file_id:
                    file_item = Files.get_file_by_id(file_id)
                    if file_item and file_item.path:
                        file_path = Storage.get_file(file_item.path)
                        with open(file_path, 'rb') as f:
                            slide_images[idx] = f.read()
        except Exception as e:
            log.warning(f'Failed to generate image for slide {idx + 1}: {e}')

    return slide_images, slide_image_urls


# ---------------------------------------------------------------------------
# Chat presentation handler
# ---------------------------------------------------------------------------


async def chat_presentation_handler(
    request: Request,
    form_data: dict,
    extra_params: dict,
    user,
    *,
    generate_chat_completion_fn=None,
    image_generations_fn=None,
    create_image_form_cls=None,
    get_last_user_message_fn=None,
    get_task_model_id_fn=None,
    add_or_update_system_message_fn=None,
):
    from open_webui.models.files import Files, FileForm
    from open_webui.storage.provider import Storage

    metadata = extra_params.get('__metadata__', {})
    chat_id = metadata.get('chat_id', None)
    __event_emitter__ = extra_params.get('__event_emitter__', None)

    if not __event_emitter__:
        return form_data

    await __event_emitter__(
        {'type': 'status', 'data': {'description': t('pres.status_designing'), 'done': False}}
    )

    _get_last = get_last_user_message_fn
    user_message = _get_last(form_data.get('messages', [])) if _get_last else ''

    models = request.app.state.MODELS
    model_id = form_data.get('model')
    task_model_id = model_id
    if get_task_model_id_fn:
        task_model_id = get_task_model_id_fn(
            model_id,
            request.app.state.config.TASK_MODEL,
            request.app.state.config.TASK_MODEL_EXTERNAL,
            models,
        )

    valid_layouts_str = ', '.join(VALID_LAYOUTS)
    slides_prompt = (
        'You are a presentation designer. '
        f'Create a JSON presentation for the following request:\n"{user_message}"\n\n'
        'Respond ONLY with a valid JSON object using this exact schema:\n'
        '{\n'
        '  "title": "Presentation Title",\n'
        '  "slides": [\n'
        '    {\n'
        f'      "layout": "<one of: {valid_layouts_str}>",\n'
        '      "title": "Slide Title",\n'
        '      "body": ["Bullet point 1", "Bullet point 2"],\n'
        '      "image_prompt": "A detailed description of the image to generate for this slide"\n'
        '    }\n'
        '  ]\n'
        '}\n\n'
        'Layout descriptions:\n'
        '- split-right: text on the left half, image on the right half. Best for content slides with supporting visuals.\n'
        '- split-left: image on the left half, text on the right half. Good for variety.\n'
        '- hero: full-bleed background image with large centered title overlay. Use for title/intro slides and dramatic impact.\n'
        '- top-strip: full-width image strip on top (~40%), title and bullets below. Good for section headers.\n\n'
        'Rules:\n'
        '- Include exactly 4 slides\n'
        '- Slide 1: use "hero" layout as a title slide\n'
        '- Slide 2-3: main content slides, vary layouts (split-right, split-left, top-strip)\n'
        '- Slide 4: conclusions/summary slide\n'
        '- Each image_prompt must be a vivid, specific description suitable for AI image generation\n'
        '- body array should have 2-4 short bullet points per slide\n'
        '- Return JSON only, no extra text or markdown'
    )

    system_message_content = ''
    try:
        if not generate_chat_completion_fn:
            raise ValueError('generate_chat_completion_fn not provided')

        result = await generate_chat_completion_fn(
            request,
            {
                'model': task_model_id,
                'messages': [{'role': 'user', 'content': slides_prompt}],
                'stream': False,
                'temperature': 0.7,
                'metadata': {
                    'task': 'presentation_generation',
                    'chat_id': chat_id,
                },
            },
            user,
        )

        raw = await drain_response_to_text(result)
        bracket_start = raw.find('{')
        bracket_end = raw.rfind('}') + 1
        if bracket_start == -1 or bracket_end <= bracket_start:
            raise ValueError('No JSON object found in LLM response')

        slides_data = json.loads(raw[bracket_start:bracket_end])
        slides = slides_data.get('slides', [])

        if len(slides) > 4:
            slides = slides[:4]
            slides_data['slides'] = slides

        for sd in slides:
            if sd.get('layout') not in VALID_LAYOUTS:
                sd['layout'] = 'split-right'
            if 'bullets' in sd and 'body' not in sd:
                sd['body'] = sd.pop('bullets')

        theme = detect_pptx_theme(user_message or '')

        # Generate images
        slide_images: dict[int, bytes] = {}
        slide_image_urls: dict[int, str] = {}
        image_gen_available = getattr(request.app.state.config, 'ENABLE_IMAGE_GENERATION', False)
        if image_gen_available:
            await __event_emitter__(
                {'type': 'status', 'data': {'description': t('pres.status_images'), 'done': False}}
            )
            slide_images, slide_image_urls = await generate_slide_images(
                request, slides, metadata, user, __event_emitter__,
                image_generations_fn=image_generations_fn,
                create_image_form_cls=create_image_form_cls,
            )

        # Build PPTX
        await __event_emitter__(
            {'type': 'status', 'data': {'description': t('pres.status_building'), 'done': False}}
        )
        file_bytes = build_pptx(slides_data, theme=theme, slide_images=slide_images)

        pres_title = slides_data.get('title', 'Presentation')
        filename_base = re.sub(r'[^\w\-]', '_', pres_title)[:50]
        filename = f'{filename_base}.pptx'
        content_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'

        file_id = str(_uuid.uuid4())
        storage_filename = f'{file_id}_{filename}'

        contents, file_path = Storage.upload_file(
            io.BytesIO(file_bytes),
            storage_filename,
            {
                'OpenWebUI-User-Email': user.email,
                'OpenWebUI-User-Id': user.id,
                'OpenWebUI-User-Name': user.name,
                'OpenWebUI-File-Id': file_id,
            },
        )

        Files.insert_new_file(
            user.id,
            FileForm(
                id=file_id,
                filename=filename,
                path=file_path,
                data={},
                meta={
                    'name': filename,
                    'content_type': content_type,
                    'size': len(file_bytes),
                },
            ),
        )

        file_url = file_id
        slide_count = len(slides)
        img_count = len(slide_images)

        await __event_emitter__(
            {'type': 'status', 'data': {'description': t('pres.status_done'), 'done': True}}
        )
        await __event_emitter__(
            {
                'type': 'files',
                'data': {
                    'files': [
                        {
                            'type': 'file',
                            'id': file_id,
                            'name': filename,
                            'url': file_url,
                            'size': len(file_bytes),
                            'content_type': content_type,
                        }
                    ]
                },
            }
        )

        slide_outline = '; '.join(
            f'{t("pres.slide_n", n=i+1)}: {s.get("title", "")}'
            for i, s in enumerate(slides)
        )
        img_note = t('pres.with_images', count=img_count) if img_count else ''
        system_message_content = t(
            'pres.system_success',
            title=pres_title,
            slide_count=slide_count,
            img_note=img_note,
            slide_outline=slide_outline,
        )

    except Exception as e:
        log.exception(e)
        await __event_emitter__(
            {'type': 'status', 'data': {'description': t('pres.status_error'), 'done': True}}
        )
        system_message_content = t('pres.system_error', error=e)

    _add_sys = add_or_update_system_message_fn
    if system_message_content and _add_sys:
        form_data['messages'] = _add_sys(system_message_content, form_data['messages'])
        messages = form_data.get('messages', [])
        for i in range(len(messages) - 1, -1, -1):
            if messages[i].get('role') == 'user':
                messages[i]['content'] = t('pres.confirm_user_message')
                break

    return form_data
